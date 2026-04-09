import re
from PIL import Image
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm


def narrative_styling(para, words, text_size, bullet=False, color=RGBColor(0, 0, 0), underline=False, line_spacing=2):
    """自定义排版样式，修复自带项目符号对齐问题"""
    for word in words:
        set_color = color
        if bullet:
            run = para.add_run()
            run.text = f"• "
            run.font.size = Pt(text_size)
            run.font.color.rgb = set_color
            bullet = False
        run = para.add_run()
        if underline:
            run.font.underline = True
        space = " " if not re.match(r'[()~]', word) else ""
        run.text = f"{word}{space}"
        run.font.size = Pt(text_size)
        run.font.color.rgb = set_color
    para.alignment = PP_ALIGN.LEFT
    para.space_after = Pt(line_spacing)
    return para


def export_to_ppt(slide_data_ls, ppt_path):
    """接收处理好的字典列表，生成并保存 PPT 文件"""
    prs = Presentation()
    slide_width_cm = 24.4
    slide_height_cm = 19.05
    prs.slide_width = Cm(slide_width_cm)
    prs.slide_height = Cm(slide_height_cm)

    for slide_data in slide_data_ls:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        key_points = slide_data.get("key_points", [])[:6]  # 最多保留6个点
        filename = slide_data.get("filename")

        # ---------------- 渲染图片 ----------------
        if filename and os.path.exists(filename):
            image = Image.open(filename)
            image_width, image_height = image.size
            image_width_cm = image_width * 0.0264583333
            image_height_cm = image_height * 0.0264583333

            while True:
                image_width_cm *= 0.9
                image_height_cm *= 0.9
                if slide_width_cm >= image_width_cm and slide_height_cm >= image_height_cm:
                    image_width_cm *= 0.9
                    image_height_cm *= 0.9
                    break

            top = Cm(2.0)
            left = Cm(13.85)
            slide.shapes.add_picture(filename, left=left, top=top, width=Cm(image_width_cm), height=Cm(image_height_cm))

        # ---------------- 渲染标题 ----------------
        txBox_title = slide.shapes.add_textbox(Cm(0), Cm(0.2), Cm(24.4), Cm(1.5))
        tx = txBox_title.text_frame
        title_p = tx.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_run = title_p.add_run()
        title_run.text = slide_data.get("topic", "无标题")
        title_run.font.name = 'Microsoft YaHei'
        title_run.font.size = Pt(22)
        title_run.font.bold = True
        tx.word_wrap = True

        # ---------------- 渲染文本要点 ----------------
        if key_points:
            # 如果没有图片，文本框可以宽一点；有图片就窄一点
            text_width = Cm(22.0) if not filename else Cm(11.0)
            txBox = slide.shapes.add_textbox(Cm(1.21), Cm(3.02), text_width, Cm(15.58))
            tf = txBox.text_frame
            tf.word_wrap = True

            point_first = tf.paragraphs[0]
            words = key_points[0].split(" ")
            narrative_styling(point_first, words, 16, bullet=True)

            for line in key_points[1:]:
                tf.add_paragraph()
                point = tf.add_paragraph()
                words = line.split(" ")
                narrative_styling(point, words, 16, bullet=True)

    prs.save(ppt_path)
    return ppt_path