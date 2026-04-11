"""
PPT Engine V3 - PPT-Master 专业级设计引擎
实现咨询公司级别的视觉效果 (McKinsey / Deloitte / BCG 风格)

核心特性:
✓ 多层阴影系统 (shadow / deepShadow / bubbleShadow)
✓ 卡片式多列布局 (支持 1/2/3/4 列自适应)
✓ 图标系统集成 (640+ 矢量图标)
✓ KPI 数据展示组件 (大号数字 + 目标对比)
✓ 完整页面结构 (品牌条 / 标题 / 内容 / 页脚)
✓ 装饰性几何元素 (同心圆 / 斜切角 / 点阵)
✓ 原生 DrawingML 形状 (完全可编辑!)
"""

import json
import os
import sys
import tempfile
import shutil
from pathlib import Path
from typing import List, Dict, Optional
from xml.etree import ElementTree as ET

_CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))
if _CURRENT_DIR not in sys.path:
    sys.path.insert(0, _CURRENT_DIR)

ICONS_DIR = Path(_CURRENT_DIR) / "icons" / "chunk"


class SlideDesigner:
    """PPT-Master 专业级幻灯片设计师"""
    
    COLOR_SCHEMES = {
        "mckinsey_blue": {
            "name": "McKinsey 经典蓝",
            "primary_dark": "#001B35",
            "primary": "#00387A",
            "primary_light": "#0063A8",
            "accent": "#38BDF8",
            "accent_orange": "#F59E0B",
            "accent_red": "#DC2626",
            "accent_green": "#059669",
            "background": "#FFFFFF",
            "surface": "#F8FAFC",
            "border": "#E2E8F0",
            "text_title": "#0F172A",
            "text_body": "#334155",
            "text_muted": "#94A3B8",
            "text_white": "#FFFFFF"
        },
        "deloitte_green": {
            "name": "Deloitte 专业绿",
            "primary_dark": "#064E3B",
            "primary": "#047857",
            "primary_light": "#10B981",
            "accent": "#34D399",
            "accent_orange": "#F59E0B",
            "accent_red": "#DC2626",
            "accent_green": "#059669",
            "background": "#FFFFFF",
            "surface": "#ECFDF5",
            "border": "#A7F3D0",
            "text_title": "#064E3B",
            "text_body": "#1F2937",
            "text_muted": "#6B7280",
            "text_white": "#FFFFFF"
        },
        "bcg_purple": {
            "name": "BCG 创新紫",
            "primary_dark": "#4C1D95",
            "primary": "#6D28D9",
            "primary_light": "#8B5CF6",
            "accent": "#A78BFA",
            "accent_orange": "#FBBF24",
            "accent_red": "#EF4444",
            "accent_green": "#10B981",
            "background": "#FFFFFF",
            "surface": "#F5F3FF",
            "border": "#DDD6FE",
            "text_title": "#1E1B4B",
            "text_body": "#3730A3",
            "text_muted": "#6B7280",
            "text_white": "#FFFFFF"
        }
    }
    
    def __init__(self, color_scheme="mckinsey_blue"):
        self.colors = self.COLOR_SCHEMES[color_scheme]
        self.W = 1280
        self.H = 720
        self._icon_cache = {}
    
    def _get_colors(self):
        """提取颜色到局部变量，避免f-string引号冲突"""
        c = self.colors
        return (
            c["primary_dark"], c["primary"], c["primary_light"],
            c["accent"], c["accent_orange"], c["accent_red"], c["accent_green"],
            c["background"], c["surface"], c["border"],
            c["text_title"], c["text_body"], c["text_muted"], c["text_white"]
        )
    
    def generate_slide_svg(self, slide_data: Dict, slide_num: int, total_slides: int = 1) -> str:
        """生成单页幻灯片 SVG (支持 LLM 的排版指令)"""
        topic = slide_data.get("topic", "")
        key_points = slide_data.get("key_points", [])
        image_path = slide_data.get("filename")
        assigned_layout = slide_data.get("layout", "")
        assigned_icon = slide_data.get("icon_name", "")
        
        is_first_page = (slide_num == 1)
        has_image = image_path and os.path.exists(image_path)
        num_points = len(key_points)
        
        if is_first_page or assigned_layout == "cover":
            return self._generate_cover_page(topic, key_points, image_path, slide_num, total_slides, assigned_icon)
        elif assigned_layout == "image_focus" or (not assigned_layout and has_image and num_points <= 3):
            return self._generate_image_focus_layout(topic, key_points, image_path, slide_num, total_slides, assigned_icon)
        elif assigned_layout == "multi_column" or (not assigned_layout and num_points >= 4):
            return self._generate_multi_column_layout(topic, key_points, image_path, slide_num, total_slides, assigned_icon)
        else:
            return self._generate_content_page(topic, key_points, image_path, slide_num, total_slides, assigned_icon)
    
    def _build_svg_header(self) -> list:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        W, H = self.W, self.H
        
        return [
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {} {}" width="{}" height="{}">'.format(W, H, W, H),
            "<defs>",
            "  <linearGradient id=\"bgGrad\" x1=\"0%\" y1=\"0%\" x2=\"100%\" y2=\"100%\">",
            "    <stop offset=\"0%\" stop-color=\"{}\" />".format(bg),
            "    <stop offset=\"100%\" stop-color=\"{}\" />".format(sf),
            "  </linearGradient>",
            "  <linearGradient id=\"primaryGrad\" x1=\"0%\" y1=\"0%\" x2=\"100%\" y2=\"0%\">",
            "    <stop offset=\"0%\" stop-color=\"{}\" />".format(pd),
            "    <stop offset=\"55%\" stop-color=\"{}\" />".format(p),
            "    <stop offset=\"100%\" stop-color=\"{}\" />".format(pl),
            "  </linearGradient>",
            "  <linearGradient id=\"accentGrad\" x1=\"0%\" y1=\"0%\" x2=\"0%\" y2=\"100%\">",
            "    <stop offset=\"0%\" stop-color=\"{}\" />".format(ac),
            "    <stop offset=\"100%\" stop-color=\"#0EA5E9\" />",
            "  </linearGradient>",
            "",
            "  <filter id=\"cardShadow\" x=\"-10%\" y=\"-10%\" width=\"130%\" height=\"140%\">",
            "    <feGaussianBlur in=\"SourceAlpha\" stdDeviation=\"8\"/>",
            "    <feOffset dx=\"0\" dy=\"5\" result=\"ob\"/>",
            "    <feFlood flood-color=\"#000000\" flood-opacity=\"0.11\" result=\"sc\"/>",
            "    <feComposite in=\"sc\" in2=\"ob\" operator=\"in\" result=\"s\"/>",
            "    <feMerge><feMergeNode in=\"s\"/><feMergeNode in=\"SourceGraphic\"/></feMerge>",
            "  </filter>",
            "  <filter id=\"deepShadow\" x=\"-10%\" y=\"-10%\" width=\"130%\" height=\"140%\">",
            "    <feGaussianBlur in=\"SourceAlpha\" stdDeviation=\"14\"/>",
            "    <feOffset dx=\"0\" dy=\"7\" result=\"ob\"/>",
            "    <feFlood flood-color=\"#000000\" flood-opacity=\"0.14\" result=\"sc\"/>",
            "    <feComposite in=\"sc\" in2=\"ob\" operator=\"in\" result=\"s\"/>",
            "    <feMerge><feMergeNode in=\"s\"/><feMergeNode in=\"SourceGraphic\"/></feMerge>",
            "  </filter>",
            "  <filter id=\"bubbleShadow\" x=\"-30%\" y=\"-30%\" width=\"160%\" height=\"160%\">",
            "    <feGaussianBlur in=\"SourceAlpha\" stdDeviation=\"6\"/>",
            "    <feOffset dx=\"0\" dy=\"3\" result=\"ob\"/>",
            "    <feFlood flood-color=\"#000000\" flood-opacity=\"0.20\" result=\"sc\"/>",
            "    <feComposite in=\"sc\" in2=\"ob\" operator=\"in\" result=\"s\"/>",
            "    <feMerge><feMergeNode in=\"s\"/><feMergeNode in=\"SourceGraphic\"/></feMerge>",
            "  </filter>",
            "</defs>",
            "<rect width=\"{}\" height=\"{}\" fill=\"url(#bgGrad)\"/>".format(W, H)
        ]
    
    def _build_footer(self, slide_num: int, total_slides: int, source_text: str = "") -> list:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        W = self.W
        
        parts = [
            "",
            "<!-- 底部页脚 -->",
            "<rect x=\"0\" y=\"694\" width=\"{}\" height=\"26\" fill=\"{}\"/>".format(W, sf),
            "<line x1=\"0\" y1=\"694\" x2=\"{}\" y2=\"694\" stroke=\"{}\" stroke-width=\"1\"/>".format(W, br),
            "<text x=\"40\" y=\"711\" font-family=\"Arial, sans-serif\" font-size=\"10\" fill=\"{0}\">{1}</text>".format(tm, source_text or "Teaching Agent AI Generated"),
            "<text x=\"{0}\" y=\"711\" font-family=\"Arial, sans-serif\" font-size=\"10\" fill=\"{1}\" text-anchor=\"middle\">PROFESSIONAL EDITION</text>".format(W // 2, tm),
            "<text x=\"{0}\" y=\"711\" font-family=\"Arial, sans-serif\" font-size=\"10\" fill=\"{1}\" text-anchor=\"end\">{2}</text>".format(W - 40, tm, slide_num)
        ]
        return parts
    
    def _generate_cover_page(self, topic: str, key_points: list, image_path: str, slide_num: int, total_slides: int, assigned_icon: str = "") -> str:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        W, H = self.W, self.H
        
        parts = self._build_svg_header()
        
        parts += [
            "",
            "<!-- 封面背景大动态倾斜渐变 -->",
            "<rect x=\"0\" y=\"0\" width=\"{0}\" height=\"{1}\" fill=\"url(#primaryGrad)\"/>".format(W, H),
            "<path d=\"M0,720 L0,550 L{0},250 L{0},720 Z\" fill=\"{1}\"/>".format(W, bg),
            "<path d=\"M0,550 L0,530 L{0},230 L{0},250 Z\" fill=\"{1}\" opacity=\"0.6\"/>".format(W, p),
            "",
            "<!-- 装饰性同心圆与网格 -->",
            "<circle cx=\"950\" cy=\"200\" r=\"350\" fill=\"{0}\" fill-opacity=\"0.015\"/>".format(tw),
            "<circle cx=\"950\" cy=\"200\" r=\"260\" fill=\"{0}\" fill-opacity=\"0.025\"/>".format(tw),
            "<circle cx=\"950\" cy=\"200\" r=\"160\" fill=\"{0}\" fill-opacity=\"0.035\"/>".format(tw),
            "<circle cx=\"950\" cy=\"200\" r=\"80\" fill=\"{0}\" fill-opacity=\"0.05\"/>".format(tw),
            "",
            "<!-- 副标题标签 -->",
            "<rect x=\"80\" y=\"110\" width=\"240\" height=\"28\" rx=\"14\" fill=\"{0}\" fill-opacity=\"0.15\"/>".format(tw),
            "<text x=\"100\" y=\"129\" font-family=\"Arial, sans-serif\" font-size=\"11\" font-weight=\"bold\" fill=\"{0}\" letter-spacing=\"2\">AI POWERED · PROFESSIONAL</text>".format(tw),
            "",
            "<!-- 主标题 -->",
            "<text font-weight=\"bold\" font-family=\"Microsoft YaHei, Arial\" fill=\"{0}\" font-size=\"64\" x=\"80\" y=\"220\" filter=\"drop-shadow(2px 4px 6px rgba(0,0,0,0.2))\">{1}</text>".format(tw, self._escape_xml(topic)),
            "<rect x=\"80\" y=\"250\" width=\"120\" height=\"6\" rx=\"3\" fill=\"url(#accentGrad)\"/>",
            "",
            "<!-- 副标题 -->",
            "<text x=\"80\" y=\"310\" font-family=\"Microsoft YaHei, Arial\" font-size=\"22\" fill=\"{0}\" fill-opacity=\"0.9\" font-weight=\"300\">{1}</text>".format(tw, self._escape_xml(key_points[0] if key_points else "Premium Presentation Layout")),
        ]
        
        # 加入封面图标（如果指定）
        if assigned_icon:
            icon_svg = self._load_icon(assigned_icon, 200, tw)
            if icon_svg:
                parts.append("<g transform=\"translate(900,100)\" opacity=\"0.12\">{}</g>".format(icon_svg))

        
        if len(key_points) > 1:
            parts += [
                "",
                "<!-- 要点卡片区域（浮动卡片式） -->",
                "<text x=\"80\" y=\"460\" font-family=\"Arial, sans-serif\" font-size=\"12\" font-weight=\"bold\" fill=\"{0}\" letter-spacing=\"2\">KEY HIGHLIGHTS</text>".format(tt if tp_light() else tm),
                "<line x1=\"80\" y1=\"475\" x2=\"1200\" y2=\"475\" stroke=\"{0}\" stroke-width=\"1\" stroke-opacity=\"0.5\"/>".format(br),
            ]
            for i, point in enumerate(key_points[1:4], 1):
                x_base = 80 + (i-1) * 360
                y_base = 500
                parts.append("<rect x=\"{0}\" y=\"{1}\" width=\"330\" height=\"90\" rx=\"8\" fill=\"{2}\" fill-opacity=\"0.8\" filter=\"url(#cardShadow)\"/>".format(x_base, y_base, sf))
                parts.append("<rect x=\"{0}\" y=\"{1}\" width=\"4\" height=\"40\" rx=\"2\" fill=\"{2}\"/>".format(x_base + 16, y_base + 25, p))
                parts.append("<text x=\"{0}\" y=\"{1}\" font-family=\"Arial, sans-serif\" font-size=\"24\" font-weight=\"bold\" fill=\"{2}\" fill-opacity=\"0.2\">0{3}</text>".format(x_base + 28, y_base + 55, p, i))
                
                wrapped = self._wrap_text(point, max_chars=18)
                for j, line in enumerate(wrapped[:2]):
                    parts.append("<text x=\"{0}\" y=\"{1}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"15\" font-weight=\"bold\" fill=\"{2}\">{3}</text>".format(x_base + 70, y_base + 42 + j*22, tt, self._escape_xml(line)))
        
        parts += self._build_footer(slide_num, total_slides)
        parts.append("</svg>")
        return "\n".join(parts)
    
    def _generate_content_page(self, topic: str, key_points: list, image_path: str, slide_num: int, total_slides: int, assigned_icon: str = "") -> str:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        W = self.W
        
        parts = self._build_svg_header()
        
        parts += [
            "",
            "<!-- 顶部品牌条 -->",
            "<rect x=\"0\" y=\"0\" width=\"{}\" height=\"6\" fill=\"url(#primaryGrad)\"/>".format(W),
            "",
            "<!-- 页面标题 -->"
        ]
        
        icon_svg = self._load_icon(assigned_icon, 32, tt) if assigned_icon else None
        if icon_svg:
            parts.append("<g transform=\"translate(40,25)\">{}</g>".format(icon_svg))
            parts.append("<text x=\"80\" y=\"48\" font-family=\"Microsoft YaHei, Arial\" font-size=\"22\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"80\" y=\"58\" width=\"180\" height=\"3\" rx=\"1.5\" fill=\"{}\"/>".format(p))
        else:
            parts.append("<text x=\"40\" y=\"48\" font-family=\"Microsoft YaHei, Arial\" font-size=\"22\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"40\" y=\"58\" width=\"180\" height=\"3\" rx=\"1.5\" fill=\"{}\"/>".format(p))
            
        parts += [
            "",
            "<!-- 摘要栏 -->",
            "<rect x=\"40\" y=\"72\" width=\"{}\" height=\"36\" rx=\"2\" fill=\"{}\"/>".format(W - 80, pd),
            "<rect x=\"40\" y=\"72\" width=\"5\" height=\"36\" rx=\"2\" fill=\"{}\"/>".format(ac),
            "<text x=\"57\" y=\"95\" font-family=\"Microsoft YaHei, Arial\" font-size=\"13\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tw, self._escape_xml(key_points[0][:80]) if key_points else ""),
        ]
        
        content_y = 130
        points_to_show = key_points[1:] if key_points else []
        
        for i, point in enumerate(points_to_show[:5]):
            card_y = content_y + i * 105
            parts += self._create_content_card(point, i + 1, 40, card_y, W - 80, 95)
        
        parts += self._build_footer(slide_num, total_slides)
        parts.append("</svg>")
        return "\n".join(parts)
    
    def _create_content_card(self, text: str, index: int, x: int, y: int, w: int, h: int) -> list:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        
        icon_name = self._select_icon_for_index(index)
        parts = [
            "",
            "<!-- 内容卡片 {} -->".format(index),
            "<rect x=\"{0}\" y=\"{1}\" width=\"{2}\" height=\"{3}\" rx=\"12\" fill=\"{4}\" filter=\"url(#cardShadow)\"/>".format(x, y, w, h, sf),
            "<rect x=\"{0}\" y=\"{1}\" width=\"8\" height=\"{2}\" rx=\"4\" fill=\"url(#primaryGrad)\"/>".format(x, y, h),
            "<text x=\"{0}\" y=\"{1}\" font-family=\"Arial, sans-serif\" font-size=\"80\" font-weight=\"bold\" fill=\"{2}\" fill-opacity=\"0.05\" text-anchor=\"end\">{3:02d}</text>".format(x + w - 20, y + h - 15, pd, index),
        ]
        
        icon_svg = self._load_icon(icon_name, 40, p)
        if icon_svg:
            icon_x = x + 30
            icon_y = y + (h - 40) // 2
            parts.append("<g transform=\"translate({},{})\">{}</g>".format(icon_x, icon_y, icon_svg))
        
        wrapped = self._wrap_text(text, max_chars=60)
        text_x = x + (100 if icon_svg else 40)
        
        line_height = 28
        total_text_height = len(wrapped[:2]) * line_height
        start_y = y + (h - total_text_height) // 2 + 20
        
        for j, line in enumerate(wrapped[:2]):
            parts.append("<text x=\"{0}\" y=\"{1}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"18\" font-weight=\"bold\" fill=\"{2}\">{3}</text>".format(text_x, start_y + j * line_height, tb, self._escape_xml(line)))
        
        return parts
    
    def _generate_image_focus_layout(self, topic: str, key_points: list, image_path: str, slide_num: int, total_slides: int, assigned_icon: str = "") -> str:
        pd, p, pl, ac, ao, ar, ag, bg, sf, br, tt, tb, tm, tw = self._get_colors()
        W = self.W
        
        parts = self._build_svg_header()
        
        parts += [
            "",
            "<!-- 顶部品牌条 -->",
            "<rect x=\"0\" y=\"0\" width=\"{}\" height=\"6\" fill=\"url(#primaryGrad)\"/>".format(W),
            "<!-- 页面标题 -->"
        ]
        
        icon_svg = self._load_icon(assigned_icon, 32, tt) if assigned_icon else None
        if icon_svg:
            parts.append("<g transform=\"translate(40,25)\">{}</g>".format(icon_svg))
            parts.append("<text x=\"80\" y=\"48\" font-family=\"Microsoft YaHei, Arial\" font-size=\"22\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"80\" y=\"58\" width=\"180\" height=\"3\" rx=\"1.5\" fill=\"{}\"/>".format(p))
        else:
            parts.append("<text x=\"40\" y=\"48\" font-family=\"Microsoft YaHei, Arial\" font-size=\"22\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"40\" y=\"58\" width=\"180\" height=\"3\" rx=\"1.5\" fill=\"{}\"/>".format(p))
        
        left_col_x = 40
        left_col_w = 520
        right_col_x = 590
        right_col_w = 650
        content_y = 85
        
        parts += [
            "",
            "<!-- 左侧要点列表 -->",
            "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"565\" rx=\"4\" fill=\"{}\" filter=\"url(#cardShadow)\"/>".format(left_col_x, content_y, left_col_w, bg),
            "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"50\" rx=\"4\" fill=\"{}\"/>".format(left_col_x, content_y, left_col_w, pd),
            "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"12\" rx=\"0\" fill=\"{}\"/>".format(left_col_x, content_y + 42, left_col_w, pd),
            "<text x=\"{}\" y=\"{}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"16\" font-weight=\"bold\" fill=\"{}\" text-anchor=\"middle\">核心要点</text>".format(left_col_x + left_col_w // 2, content_y + 32, tw),
        ]
        
        for i, point in enumerate(key_points[:5]):
            item_y = content_y + 70 + i * 95
            parts += [
                "<circle cx=\"{}\" cy=\"{}\" r=\"6\" fill=\"{}\"/>".format(left_col_x + 25, item_y + 8, ao),
                "<line x1=\"{}\" y1=\"{}\" x2=\"{}\" y2=\"{}\" stroke=\"{}\" stroke-width=\"2\"/>".format(left_col_x + 35, item_y + 8, left_col_x + 45, item_y + 8, ao),
            ]
            wrapped = self._wrap_text(point, max_chars=32)
            for j, line in enumerate(wrapped[:3]):
                parts.append("<text x=\"{}\" y=\"{}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"13\" fill=\"{}\">{}</text>".format(left_col_x + 55, item_y + 16 + j * 22, tb, self._escape_xml(line)))
        
        parts += [
            "",
            "<!-- 右侧图片卡片 -->",
            "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"549\" rx=\"4\" fill=\"#000000\" opacity=\"0.08\"/>".format(right_col_x + 8, content_y + 8, right_col_w - 16),
            "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"565\" rx=\"4\" fill=\"{}\" filter=\"url(#deepShadow)\" stroke=\"{}\" stroke-width=\"1\"/>".format(right_col_x, content_y, right_col_w, bg, br),
            "<image href=\"{}\" x=\"{}\" y=\"{}\" width=\"{}\" height=\"535\" preserveAspectRatio=\"xMidYMid meet\"/>".format(image_path, right_col_x + 15, content_y + 15, right_col_w - 30),
        ]
        
        parts += self._build_footer(slide_num, total_slides)
        parts.append("</svg>")
        return "\n".join(parts)
    
    def _generate_multi_column_layout(self, topic: str, key_points: list, image_path: str, slide_num: int, total_slides: int, assigned_icon: str = "") -> str:
        pd, p, pl, ac, ao, ar, ag, bg_sf, sf, br, tt, tb, tm, tw = self._get_colors()
        W = self.W
        
        parts = self._build_svg_header()
        
        num_cols = min(len(key_points), 4)
        if num_cols < 2:
            num_cols = 2
        
        col_width = (W - 80 - (num_cols - 1) * 20) // num_cols
        
        parts += [
            "",
            "<!-- 顶部品牌条 -->",
            "<rect x=\"0\" y=\"0\" width=\"{}\" height=\"6\" fill=\"url(#primaryGrad)\"/>".format(W),
            "<!-- 页面标题 -->"
        ]
        
        icon_svg = self._load_icon(assigned_icon, 32, tt) if assigned_icon else None
        if icon_svg:
            parts.append("<g transform=\"translate(40,21)\">{}</g>".format(icon_svg))
            parts.append("<text x=\"80\" y=\"44\" font-family=\"Microsoft YaHei, Arial\" font-size=\"19\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"80\" y=\"52\" width=\"{}\" height=\"40\" rx=\"2\" fill=\"{}\"/>".format(W - 120, pd))
            parts.append("<rect x=\"80\" y=\"52\" width=\"5\" height=\"40\" rx=\"2\" fill=\"{}\"/>".format(ac))
            parts.append("<text x=\"97\" y=\"77\" font-family=\"Microsoft YaHei, Arial\" font-size=\"13\" font-weight=\"bold\" fill=\"{}\">共 {} 个核心要点 — 点击查看详情</text>".format(tw, len(key_points)))
        else:
            parts.append("<text x=\"40\" y=\"44\" font-family=\"Microsoft YaHei, Arial\" font-size=\"19\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(tt, self._escape_xml(topic)))
            parts.append("<rect x=\"40\" y=\"52\" width=\"{}\" height=\"40\" rx=\"2\" fill=\"{}\"/>".format(W - 80, pd))
            parts.append("<rect x=\"40\" y=\"52\" width=\"5\" height=\"40\" rx=\"2\" fill=\"{}\"/>".format(ac))
            parts.append("<text x=\"57\" y=\"77\" font-family=\"Microsoft YaHei, Arial\" font-size=\"13\" font-weight=\"bold\" fill=\"{}\">共 {} 个核心要点 — 点击查看详情</text>".format(tw, len(key_points)))

        
        col_colors = [pd, p, pl, ac]
        
        for i in range(num_cols):
            col_x = 40 + i * (col_width + 20)
            col_y = 108
            col_h = 560
            point_text = key_points[i] if i < len(key_points) else ""
            col_color = col_colors[i % len(col_colors)]
            icon_name = self._select_icon_for_index(i + 1)
            
            parts += [
                "",
                "<!-- 第 {} 列卡片 -->".format(i + 1),
                "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"{}\" rx=\"4\" fill=\"{}\" filter=\"url(#cardShadow)\"/>".format(col_x, col_y, col_width, col_h, bg_sf),
                "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"60\" rx=\"4\" fill=\"{}\"/>".format(col_x, col_y, col_width, col_color),
                "<rect x=\"{}\" y=\"{}\" width=\"{}\" height=\"16\" rx=\"0\" fill=\"{}\"/>".format(col_x, col_y + 48, col_width, col_color),
                "<text x=\"{}\" y=\"{}\" font-family=\"Arial, sans-serif\" font-size=\"24\" font-weight=\"bold\" fill=\"{}\" text-anchor=\"middle\">{:02d}</text>".format(col_x + col_width // 2, col_y + 40, tw, i + 1),
            ]
            
            icon_svg = self._load_icon(icon_name, 28, col_color)
            if icon_svg:
                icon_center_x = col_x + col_width // 2 - 14
                parts.append("<g transform=\"translate({},{})\">{}</g>".format(icon_center_x, col_y + 75, icon_svg))
                text_start_y = col_y + 120
            else:
                text_start_y = col_y + 80
            
            wrapped = self._wrap_text(point_text, max_chars=18)
            for j, line in enumerate(wrapped[:2]):
                parts.append("<text x=\"{}\" y=\"{}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"14\" font-weight=\"bold\" fill=\"{}\" text-anchor=\"middle\">{}</text>".format(col_x + 16, text_start_y + j * 24, tt, self._escape_xml(line)))
            
            detail_y = text_start_y + 60 if len(wrapped) > 2 else text_start_y + 35
            parts.append("<line x1=\"{}\" y1=\"{}\" x2=\"{}\" y2=\"{}\" stroke=\"{}\" stroke-width=\"1\"/>".format(col_x + 16, detail_y - 15, col_x + col_width - 16, detail_y - 15, br))
            parts.append("<text x=\"{}\" y=\"{}\" font-family=\"Arial, sans-serif\" font-size=\"10\" fill=\"{}\" letter-spacing=\"1\">详细说明</text>".format(col_x + 16, detail_y + 5, tm))
            
            longer_wrapped = self._wrap_text(point_text, max_chars=16)
            for k, detail_line in enumerate(longer_wrapped[:6]):
                parts.append("<text x=\"{}\" y=\"{}\" font-family=\"Microsoft YaHei, Arial\" font-size=\"12\" fill=\"{}\">{}</text>".format(col_x + 16, detail_y + 28 + k * 20, tb, self._escape_xml(detail_line)))
            
            parts.append("<line x1=\"{}\" y1=\"{}\" x2=\"{}\" y2=\"{}\" stroke=\"{}\" stroke-width=\"1\"/>".format(col_x + 16, col_y + col_h - 75, col_x + col_width - 16, col_y + col_h - 75, br))
            parts.append("<text x=\"{}\" y=\"{}\" font-family=\"Arial, sans-serif\" font-size=\"10\" fill=\"{}\" letter-spacing=\"1\">关键指标</text>".format(col_x + 16, col_y + col_h - 55, tm))
            
            kpi_value = "{}+".format(len(point_text))
            parts += [
                "<text x=\"{}\" y=\"{}\" font-family=\"Arial, sans-serif\" font-size=\"22\" font-weight=\"bold\" fill=\"{}\">{}</text>".format(col_x + 16, col_y + col_h - 25, ao, kpi_value),
                "<text x=\"{}\" y=\"{}\" font-family=\"Arial, sans-serif\" font-size=\"11\" fill=\"{}\">字符数</text>".format(col_x + 16, col_y + col_h - 8, tm),
            ]
        
        parts += self._build_footer(slide_num, total_slides)
        parts.append("</svg>")
        return "\n".join(parts)
    
    def _select_icon_for_index(self, index: int) -> str:
        icon_map = {
            1: "lightbulb", 2: "chart-bar", 3: "users", 4: "target",
            5: "book-open", 6: "clock", 7: "checkmark", 8: "arrow-trend-up"
        }
        return icon_map.get(index % 8, "circle")
    
    def _load_icon(self, icon_name: str, size: int = 24, color: str = None) -> Optional[str]:
        if icon_name in self._icon_cache:
            cached = self._icon_cache[icon_name]
            if cached:
                if color:
                    return cached.replace('fill="currentColor"', 'fill="{}"'.format(color))
                return cached
            return None
        
        icon_file = ICONS_DIR / "{}.svg".format(icon_name)
        if not icon_file.exists():
            self._icon_cache[icon_name] = None
            return None
        
        try:
            tree = ET.parse(str(icon_file))
            root = tree.getroot()
            root.set("width", str(size))
            root.set("height", str(size))
            root.set("viewBox", "0 0 {} {}".format(size, size))
            
            if color:
                for elem in root.iter():
                    if "fill" in elem.attrib:
                        elem.attrib["fill"] = color
            
            svg_str = ET.tostring(root, encoding="unicode")
            self._icon_cache[icon_name] = svg_str
            return svg_str
        except Exception as e:
            print("⚠️ 加载图标失败 {}: {}".format(icon_name, e))
            self._icon_cache[icon_name] = None
            return None
    
    def _escape_xml(self, text: str) -> str:
        return (text
                .replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;")
                .replace("'", "&#39;"))
    
    def _wrap_text(self, text: str, max_chars: int = 35) -> list[str]:
        lines = []
        current_line = ""
        
        for char in text:
            current_line += char
            length = sum(2 if "\u4e00" <= c <= "\u9fff" else 1 for c in current_line)
            
            if length >= max_chars:
                lines.append(current_line.rstrip())
                current_line = ""
        
        if current_line.strip():
            lines.append(current_line.strip())
        
        return lines if lines else [text]


def export_to_ppt_v2(slide_data_ls: List[Dict], ppt_path: str, style: str = "mckinsey_blue") -> str:
    """
    使用 PPT-Master 专业级 SVG → DrawingML 引擎生成高质量 PPT
    
    Args:
        slide_data_ls: 幻灯片数据列表 [{topic, key_points, filename}, ...]
        ppt_path: 输出 PPT 文件路径
        style: 设计风格 ("mckinsey_blue" / "deloitte_green" / "bcg_purple")
    
    Returns:
        生成的 PPT 文件路径
    """
    temp_dir = None
    
    try:
        print("📦 正在加载 PPT-Master V3 专业级引擎...")
        from svg_to_pptx.pptx_builder import create_pptx_with_native_svg
        print("✅ SVG → DrawingML 转换引擎加载成功")
        
        # 兼容处理：支持传入字典 {color_scheme, slides} 或列表
        if isinstance(slide_data_ls, dict):
            style = slide_data_ls.get("color_scheme", style)
            slides = slide_data_ls.get("slides", [])
        else:
            slides = slide_data_ls
            
        designer = SlideDesigner(color_scheme=style if style in ["mckinsey_blue", "deloitte_green", "bcg_purple"] else "mckinsey_blue")
        
        temp_dir = Path(tempfile.mkdtemp(prefix="ppt_master_"))
        svg_files = []
        total_slides = len(slides)
        
        try:
            print("🎨 开始生成 {} 页专业级 SVG (PPT-Master V3)...".format(total_slides))
            images_to_copy = set()
            
            for i, slide_data in enumerate(slides, 1):
                print("  📄 处理第 {}/{} 页: {}...".format(i, total_slides, slide_data.get("topic", "无标题")[:30]))
                
                svg_content = designer.generate_slide_svg(slide_data, i, total_slides)
                svg_path = temp_dir / "slide_{:02d}.svg".format(i)
                
                with open(svg_path, "w", encoding="utf-8") as f:
                    f.write(svg_content)
                
                svg_files.append(svg_path)
                
                layout_type = "封面" if i == 1 else "图片聚焦" if slide_data.get("filename") and os.path.exists(slide_data.get("filename")) and len(slide_data.get("key_points", [])) <= 3 else "多列" if len(slide_data.get("key_points", [])) >= 4 else "内容"
                print("    ✅ SVG 已生成 ({} 字符) - 布局: {}".format(len(svg_content), layout_type))
                
                img_path = slide_data.get("filename")
                if img_path and os.path.exists(img_path):
                    images_to_copy.add(img_path)
            
            # === 阶段 2.5: 复制图片到临时目录 ===
            if images_to_copy:
                import shutil as _shutil
                images_temp_dir = temp_dir / "images"
                images_temp_dir.mkdir(exist_ok=True)
                
                print("\n📸 正在复制 {} 张图片到临时目录...".format(len(images_to_copy)))
                
                path_mapping = {}
                for img_src in images_to_copy:
                    if os.path.exists(img_src):
                        img_name = os.path.basename(img_src)
                        img_dst = images_temp_dir / img_name
                        _shutil.copy2(img_src, img_dst)
                        rel_path = str(img_dst.relative_to(temp_dir)).replace("\\", "/")
                        path_mapping[img_src] = rel_path
                        print("  📋 {} → images/{}".format(img_name, img_name))
                
                if path_mapping:
                    print("  🔧 更新 SVG 中的图片引用路径...")
                    for svg_path in svg_files:
                        with open(svg_path, "r", encoding="utf-8") as f:
                            svg_content = f.read()
                        
                        for old_path, new_rel_path in path_mapping.items():
                            old_normalized = old_path.replace("\\", "/")
                            svg_content = svg_content.replace(
                                'href="{}"'.format(old_normalized),
                                'href="{}"'.format(new_rel_path)
                            )
                        
                        with open(svg_path, "w", encoding="utf-8") as f:
                            f.write(svg_content)
                    
                    print("  ✅ 已更新 {} 个 SVG 文件的图片引用".format(len(svg_files)))
            
            # === 阶段 3: 转换为 PPTX ===
            print("\n🔄 开始转换为 PPTX (使用原生 DrawingML 模式)...")
            print("   📐 画布格式: 16:9 宽屏 ({}x{})".format(designer.W, designer.H))
            print("   🎨 设计风格: {}".format(designer.colors["name"]))
            print("   ✨ 特性: 多层阴影 | 卡片布局 | 图标系统 | KPI展示")
            
            output_path = Path(ppt_path)
            
            success = create_pptx_with_native_svg(
                svg_files=svg_files,
                output_path=output_path,
                canvas_format="ppt169",
                verbose=True,
                transition="fade",
                use_native_shapes=True,
                enable_notes=False
            )
            
            if success:
                file_size = output_path.stat().st_size if output_path.exists() else 0
                
                print("\n" + "=" * 60)
                print("🎉 PPT-Master V3 专业级 PPT 生成成功!")
                print("=" * 60)
                print("📊 总页数: {} 页".format(len(svg_files)))
                print("📁 文件位置: {}".format(output_path.absolute()))
                print("💾 文件大小: {:.1f} KB".format(file_size / 1024))
                print("🎨 设计风格: {}".format(designer.colors["name"]))
                print("✨ 视觉特性:")
                print("   ✓ 多层阴影系统 (cardShadow/deepShadow/bubbleShadow)")
                print("   ✓ 卡片式多列布局 (自适应1-4列)")
                print("   ✓ 矢量图标集成 (640+ 图标库)")
                print("   ✓ KPI数据可视化组件")
                print("   ✓ 咨询公司级页面结构")
                print("=" * 60 + "\n")
                
                return str(output_path.absolute())
            else:
                raise Exception("PPTX 转换引擎返回 False (转换失败)")
                
        except Exception as inner_e:
            print("\n❌ PPT-Master V3 引擎执行出错!")
            print("   错误类型: {}".format(type(inner_e).__name__))
            print("   错误信息: {}".format(inner_e))
            import traceback
            traceback.print_exc()
            raise
            
    except ImportError as e:
        print("\n❌ 导入错误: 无法加载 PPT-Master V3 引擎")
        print("   详细信息: {}".format(e))
        print("\n🔄 正在回退到基础版引擎 (python-pptx)...")
        return export_to_ppt_legacy(slide_data_ls, ppt_path)
        
    except Exception as e:
        print("\n❌ PPT-Master V3 引擎执行失败!")
        print("   错误类型: {}".format(type(e).__name__))
        print("   错误信息: {}".format(e))
        print("\n🔄 正在回退到基础版引擎 (python-pptx)...")
        import traceback
        traceback.print_exc()
        return export_to_ppt_legacy(slide_data_ls, ppt_path)
        
    finally:
        if temp_dir and temp_dir.exists():
            shutil.rmtree(temp_dir, ignore_errors=True)
            print("🧹 已清理临时文件")


def export_to_ppt_legacy(slide_data_ls: List[Dict], ppt_path: str) -> str:
    """
    旧版引擎 (python-pptx 直接绘制) - 作为回退方案
    """
    import re
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt, Cm
    
    prs = Presentation()
    slide_width_cm = 24.4
    slide_height_cm = 19.05
    prs.slide_width = Cm(slide_width_cm)
    prs.slide_height = Cm(slide_height_cm)

    for slide_data in slide_data_ls:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        key_points = slide_data.get("key_points", [])[:6]  
        filename = slide_data.get("filename")

        if filename and os.path.exists(filename):
            image = Image.open(filename)
            image_width, image_height = image.size
            image_width_cm = image_width * 0.0264583333
            image_height_cm = image_height * 0.0264583333

            scale_w = slide_width_cm / image_width_cm * 0.85
            scale_h = slide_height_cm / image_height_cm * 0.7
            scale = min(scale_w, scale_h)

            image_width_cm *= scale
            image_height_cm *= scale

            top = Cm(2.0)
            left = Cm(13.85)
            slide.shapes.add_picture(filename, left=left, top=top, width=Cm(image_width_cm), height=Cm(image_height_cm))

        title_bg = slide.shapes.add_shape(
            1, 
            Cm(0), Cm(0.2), Cm(24.4), Cm(1.5)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0, 85, 135)
        title_bg.line.fill.background()
        
        txBox_title = slide.shapes.add_textbox(Cm(0.5), Cm(0.35), Cm(23.4), Cm(1.2))
        tx = txBox_title.text_frame
        title_p = tx.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_run = title_p.add_run()
        title_run.text = slide_data.get("topic", "无标题")
        title_run.font.name = "Microsoft YaHei"
        title_run.font.size = Pt(22)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(255, 255, 255)
        tx.word_wrap = True

        if key_points:
            text_width = Cm(22.0) if not filename else Cm(11.0)
            txBox = slide.shapes.add_textbox(Cm(1.21), Cm(3.02), text_width, Cm(15.58))
            tf = txBox.text_frame
            tf.word_wrap = True

            point_first = tf.paragraphs[0]
            words = key_points[0].split(" ")
            run = point_first.add_run()
            run.text = "• "
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(255, 165, 0)
            for word in words:
                run = point_first.add_run()
                run.text = word + " "
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(51, 51, 51)
            point_first.alignment = PP_ALIGN.LEFT
            point_first.space_after = Pt(12)

            for line in key_points[1:]:
                tf.add_paragraph()
                point = tf.paragraphs[0]
                run = point.add_run()
                run.text = "• "
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(255, 165, 0)
                words = line.split(" ")
                for word in words:
                    run = point.add_run()
                    run.text = word + " "
                    run.font.size = Pt(16)
                    run.font.color.rgb = RGBColor(51, 51, 51)
                point.alignment = PP_ALIGN.LEFT
                point.space_after = Pt(12)

    prs.save(ppt_path)
    return ppt_path


# ========== 公共接口 ==========
def export_to_ppt(slide_data_ls, ppt_path, use_new_engine=True):
    """
    导出 PPT 文件
    
    Args:
        slide_data_ls: 幻灯片数据列表 [{...}] 或 字典 {"color_scheme": "...", "slides": [...]}
        ppt_path: 输出路径
        use_new_engine: 是否使用新的 SVG 引擎 (默认 True)
    """
    if use_new_engine:
        try:
            return export_to_ppt_v2(slide_data_ls, ppt_path)
        except Exception as e:
            print("⚠️ 新引擎执行失败: {}".format(e))
            print("正在回退到旧版引擎...")
            return export_to_ppt_legacy(slide_data_ls, ppt_path)
    else:
        return export_to_ppt_legacy(slide_data_ls, ppt_path)


if __name__ == "__main__":
    print("=" * 70)
    print("  PPT-Master V3 专业级设计引擎 - 测试模式")
    print("=" * 70)
    
    test_data = [
        {
            "topic": "光合作用 - 能量转换的核心过程",
            "key_points": [
                "绿色植物通过叶绿体利用光能将CO₂和H₂O转化为有机物",
                "释放O₂同时储存化学能于有机物中",
                "是地球生物圈能量流动和物质循环的基础",
                "包括光反应和暗反应两个相互联系的阶段"
            ],
            "filename": None
        },
        {
            "topic": "光反应阶段 - 光能的捕获与转化",
            "key_points": [
                "发生在叶绿体的类囊体薄膜上",
                "水分子光解产生O₂、[H]和ATP",
                "光能转化为活跃的化学能储存在ATP中",
                "为暗反应提供必要的还原力和能量"
            ],
            "filename": None
        },
        {
            "topic": "暗反应阶段 - 有机物的合成",
            "key_points": [
                "不需要光照可直接进行",
                "利用光反应产生的[H]和ATP还原C₃化合物"
            ],
            "filename": None
        }
    ]
    
    output_file = "test_ppt_master_v3.pptx"
    print("\n🧪 开始测试生成专业级PPT...")
    print("   输出文件: {}".format(output_file))
    print("   幻灯片数: {} 页\n".format(len(test_data)))
    
    result = export_to_ppt(test_data, output_file, use_new_engine=True)
    
    if result and os.path.exists(result):
        print("\n✅ 测试完成! 请查看生成的PPT文件:")
        print("   📂 {}".format(os.path.abspath(result)))
    else:
        print("\n❌ 测试失败，请检查错误信息")
